"""
backend.kb_doc_parsers — 各类文档抽文本。

依赖库需在服务器安装(见 requirements 注释)。每个解析器对缺库做友好报错,
不让整个服务崩。OCR 通过注入的 ocr_fn(默认 None=不启用),由管理后台开关控制。

verified: 编排/路由/报告 已离线自测;真正的 pdf/docx/xlsx/pptx 抽取在装库的服务器上运行。
"""
from __future__ import annotations

import io
from typing import Callable, Optional

# OCR 注入点:bytes(图片/页面)-> str。生产传 PaddleOCR 封装;None=不启用 OCR。
OcrFn = Optional[Callable[[bytes], str]]


class ParserError(Exception):
    pass


def _need(lib: str):
    raise ParserError(f"需要在服务器安装依赖库: {lib}(本功能由龙虾部署时 pip 安装)")


def parse_text(data: bytes, **_) -> str:
    return data.decode("utf-8", errors="replace")


def parse_csv(data: bytes, **_) -> str:
    import csv
    text = data.decode("utf-8", errors="replace")
    rows = list(csv.reader(io.StringIO(text)))
    return "\n".join("\t".join(r) for r in rows)


def parse_docx(data: bytes, **_) -> str:
    try:
        import docx  # python-docx
    except ImportError:
        _need("python-docx")
    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def parse_xlsx(data: bytes, **_) -> str:
    try:
        import openpyxl
    except ImportError:
        _need("openpyxl")
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# 工作表: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                out.append("\t".join(cells))
    return "\n".join(out)


def parse_pptx(data: bytes, **_) -> str:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        _need("python-pptx")
    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# 幻灯片 {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        out.append(t)
    return "\n".join(out)


def parse_pdf(data: bytes, ocr_fn: OcrFn = None, **_) -> str:
    """优先抽文本层;无文本层且启用 OCR → 逐页 OCR。"""
    try:
        import pdfplumber
    except ImportError:
        _need("pdfplumber")
    out: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            txt = page.extract_text() or ""
            if txt.strip():
                out.append(txt)
            elif ocr_fn:                       # 扫描页:渲染成图再 OCR
                img = page.to_image(resolution=200).original
                buf = io.BytesIO(); img.save(buf, format="PNG")
                out.append(ocr_fn(buf.getvalue()))
    return "\n".join(out)


def parse_ocr_image(data: bytes, ocr_fn: OcrFn = None, **_) -> str:
    if not ocr_fn:
        raise ParserError("图片需启用 OCR(管理后台开启 OCR 开关)")
    return ocr_fn(data)


# 解析器名 → 函数
PARSERS: dict[str, Callable[..., str]] = {
    "text": parse_text, "csv": parse_csv, "docx": parse_docx,
    "xlsx": parse_xlsx, "pptx": parse_pptx, "pdf": parse_pdf,
    "ocr_image": parse_ocr_image,
}


def get_parser(name: str) -> Callable[..., str]:
    fn = PARSERS.get(name)
    if not fn:
        raise ParserError(f"未知解析器: {name}")
    return fn
